"""
AetherAI — Document Agent  (Stage 4 — hardened, patched)

Fix applied
───────────
FIX 9  PPTX slides were getting irrelevant stock photos (Golden Gate Bridge
       for a Roman Empire presentation) because image keywords were constructed
       by naively joining the slide title with the topic, producing long phrases
       that Unsplash / Picsum couldn't interpret meaningfully.

       New approach:
         • _build_image_query() strips English stop-words, de-duplicates tokens,
           and keeps the 4 most meaningful (longest) content words.
         • fetch_images_for_slides() now uses per-slide keyword construction
           so each slide gets a targeted query rather than a generic topic dump.
         • The Unsplash URL format is kept the same; keyword quality is the fix.
"""

import asyncio
import logging
import re
import json
import random
from datetime import datetime
from io import BytesIO
from pathlib import Path
from typing import Optional
from urllib.parse import quote

import httpx

from agents import BaseAgent

logger = logging.getLogger(__name__)

OUTPUT_DIR = Path(__file__).parent.parent.parent / "output"
OUTPUT_DIR.mkdir(exist_ok=True)

# ── Image sources ─────────────────────────────────────────────────────────────
#
# Source priority:
#   1. Wikimedia REST API  — returns the actual Wikipedia article image for a
#      topic/keyword. Highly accurate, free, no key, no rate limit for small use.
#      URL: https://en.wikipedia.org/api/rest_v1/page/summary/{keyword}
#           → response.thumbnail.source
#
#   2. loremflickr         — Flickr photos tagged with the keyword. Decent relevance,
#      occasionally returns cached/repeated images for close queries.
#
#   3. Picsum              — Fully random. Absolute last resort.
#
WIKIMEDIA_URL   = "https://en.wikipedia.org/api/rest_v1/page/summary/{keyword}"
LOREMFLICKR_URL = "https://loremflickr.com/900/600/{query}?lock={seed}"
PICSUM_URL      = "https://picsum.photos/seed/{seed}/900/600"

HEADERS = {
    "User-Agent": (
        "AetherAI-DocumentAgent/1.0 (https://github.com/pat1322/aetherAI; "
        "contact via GitHub) python-httpx"
    )
}

_IMAGE_SEM = asyncio.Semaphore(4)  # up from 3 — all slides now fetch images

# FIX 9: Stop-words to strip before building image search queries
_STOP_WORDS = {
    "the", "a", "an", "and", "or", "of", "in", "on", "at", "to", "for",
    "with", "by", "is", "are", "was", "were", "be", "been", "being",
    "have", "has", "had", "do", "does", "did", "will", "would", "could",
    "should", "may", "might", "shall", "can", "its", "their", "our",
    # Generic slide filler words that add no image-search value
    "overview", "introduction", "conclusion", "summary", "impact",
    "effects", "analysis", "history", "about", "what", "how", "why",
    "when", "key", "main", "top", "role", "rise", "fall", "age",
    "era", "period", "chapter", "section", "part", "slide",
}


def _build_image_query(slide_title: str, topic: str) -> str:
    """
    Construct a focused image search query from a slide title and the overall
    topic. Stop-words are removed, tokens are de-duplicated, and the 4 most
    descriptive (longest) words are selected so Unsplash gets a specific,
    meaningful query rather than a generic phrase dump.
    """
    combined = f"{slide_title} {topic}".lower()
    # Extract word tokens (3+ chars, letters only)
    words = re.findall(r'\b[a-z]{3,}\b', combined)
    # Remove stop-words
    filtered = [w for w in words if w not in _STOP_WORDS]
    # De-duplicate while preserving order
    seen: set[str] = set()
    unique: list[str] = []
    for w in filtered:
        if w not in seen:
            seen.add(w)
            unique.append(w)
    # Pick the 4 most specific (longest) words for a sharper query
    unique.sort(key=len, reverse=True)
    query = " ".join(unique[:4])
    return query or slide_title[:30]


# ── Themes ────────────────────────────────────────────────────────────────────

THEMES = [
    {
        "name": "Midnight Executive",
        "bg_dark":    (0x1E, 0x27, 0x61),
        "bg_light":   (0xF4, 0xF6, 0xFF),
        "accent":     (0x4A, 0x90, 0xD9),
        "text_dark":  (0x1E, 0x27, 0x61),
        "text_light": (0xFF, 0xFF, 0xFF),
        "muted":      (0xCA, 0xDC, 0xFC),
        "alt_row":    (0xE8, 0xEE, 0xFD),
        "header_font": "Georgia",
        "body_font":   "Calibri",
    },
    {
        "name": "Coral Energy",
        "bg_dark":    (0x2F, 0x3C, 0x7E),
        "bg_light":   (0xFF, 0xFB, 0xF0),
        "accent":     (0xF9, 0x61, 0x67),
        "text_dark":  (0x2F, 0x3C, 0x7E),
        "text_light": (0xFF, 0xFF, 0xFF),
        "muted":      (0xF9, 0xE7, 0x95),
        "alt_row":    (0xFF, 0xF0, 0xF0),
        "header_font": "Arial Black",
        "body_font":   "Arial",
    },
    {
        "name": "Ocean Gradient",
        "bg_dark":    (0x06, 0x5A, 0x82),
        "bg_light":   (0xF0, 0xF7, 0xFB),
        "accent":     (0x02, 0xC3, 0x9A),
        "text_dark":  (0x06, 0x5A, 0x82),
        "text_light": (0xFF, 0xFF, 0xFF),
        "muted":      (0x9B, 0xC8, 0xDB),
        "alt_row":    (0xE3, 0xF4, 0xF8),
        "header_font": "Trebuchet MS",
        "body_font":   "Calibri",
    },
    {
        "name": "Warm Terracotta",
        "bg_dark":    (0xB8, 0x50, 0x42),
        "bg_light":   (0xF5, 0xF3, 0xEE),
        "accent":     (0xA7, 0xBE, 0xAE),
        "text_dark":  (0x4A, 0x1A, 0x14),
        "text_light": (0xFF, 0xFF, 0xFF),
        "muted":      (0xE7, 0xE8, 0xD1),
        "alt_row":    (0xF9, 0xF1, 0xEC),
        "header_font": "Cambria",
        "body_font":   "Calibri",
    },
    {
        "name": "Cherry Bold",
        "bg_dark":    (0x99, 0x00, 0x11),
        "bg_light":   (0xFC, 0xF6, 0xF5),
        "accent":     (0x2F, 0x3C, 0x7E),
        "text_dark":  (0x33, 0x00, 0x00),
        "text_light": (0xFF, 0xFF, 0xFF),
        "muted":      (0xCC, 0x88, 0x88),
        "alt_row":    (0xFD, 0xED, 0xED),
        "header_font": "Impact",
        "body_font":   "Arial",
    },
    {
        "name": "Teal Trust",
        "bg_dark":    (0x02, 0x80, 0x90),
        "bg_light":   (0xF0, 0xFB, 0xFC),
        "accent":     (0x02, 0xC3, 0x9A),
        "text_dark":  (0x01, 0x3A, 0x40),
        "text_light": (0xFF, 0xFF, 0xFF),
        "muted":      (0x7F, 0xC8, 0xCF),
        "alt_row":    (0xE2, 0xF7, 0xF8),
        "header_font": "Calibri",
        "body_font":   "Calibri Light",
    },
]


def pick_theme() -> dict:
    return random.choice(THEMES)


# ── Image helpers ─────────────────────────────────────────────────────────────

async def fetch_image_bytes(keyword: str, slide_index: int, timeout: float = 10.0) -> Optional[bytes]:
    """
    Fetch a topic-relevant image for a slide.
    Tries Wikimedia (Wikipedia article thumbnail) first — most accurate.
    Falls back to loremflickr, then Picsum.
    Each source uses a different keyword form for best results.
    """
    async with _IMAGE_SEM:
        words = keyword.split()
        seed  = abs(hash(keyword)) % 9000 + slide_index * 17

        # ── Attempt 1: Wikimedia REST API ──────────────────────────────────────
        # Try each word combo from most-specific to least until we get a thumbnail.
        # Wikipedia's summary endpoint returns the article's main image, which is
        # always directly relevant to the search term.
        wiki_candidates = []
        if len(words) >= 2:
            wiki_candidates.append("_".join(words[:2]))   # e.g. "solar_system"
        wiki_candidates.append(words[0])                   # e.g. "solar"
        if len(words) >= 3:
            wiki_candidates.append("_".join(words[:3]))   # e.g. "solar_system_planets"

        for wkw in wiki_candidates:
            try:
                wkw_encoded = quote(wkw.replace(" ", "_"), safe="")
                wiki_url = WIKIMEDIA_URL.format(keyword=wkw_encoded)
                async with httpx.AsyncClient(
                    headers={"User-Agent": HEADERS["User-Agent"]},
                    follow_redirects=True, timeout=timeout
                ) as c:
                    r = await c.get(wiki_url)
                    if r.status_code == 200:
                        data = r.json()
                        thumb = data.get("thumbnail", {}).get("source", "")
                        if thumb:
                            # Download the actual image from the thumbnail URL
                            r2 = await c.get(thumb)
                            if r2.status_code == 200 and r2.headers.get("content-type","").startswith("image"):
                                logger.info(f"[DocumentAgent] Wikimedia OK slide {slide_index} kw='{wkw}'")
                                return r2.content
            except Exception as e:
                logger.debug(f"[DocumentAgent] Wikimedia failed '{wkw}': {e}")

        # ── Attempt 2: loremflickr ─────────────────────────────────────────────
        try:
            query = ",".join(words[:4]) or keyword[:40]
            url   = LOREMFLICKR_URL.format(query=query, seed=seed)
            async with httpx.AsyncClient(
                headers={"User-Agent": "Mozilla/5.0"},
                follow_redirects=True, timeout=timeout
            ) as c:
                r  = await c.get(url)
                ct = r.headers.get("content-type", "")
                if r.status_code == 200 and ct.startswith("image"):
                    logger.info(f"[DocumentAgent] loremflickr OK slide {slide_index} query='{query}'")
                    return r.content
        except Exception as e:
            logger.debug(f"[DocumentAgent] loremflickr failed: {e}")

        # ── Attempt 3: Picsum (random — last resort) ───────────────────────────
        try:
            async with httpx.AsyncClient(follow_redirects=True, timeout=timeout) as c:
                r = await c.get(PICSUM_URL.format(seed=seed))
                if r.status_code == 200:
                    logger.info(f"[DocumentAgent] Picsum fallback OK slide {slide_index}")
                    return r.content
        except Exception as e:
            logger.debug(f"[DocumentAgent] Picsum failed: {e}")

        return None


async def fetch_images_for_slides(slides: list, topic: str) -> dict[int, bytes]:
    """
    Fetch a relevant image for EVERY content slide.
    Each slide gets its own focused keyword from its title + the overall topic.
    Removed the old `if i % 2 == 0` filter that caused alternate slides to be text-only.
    """
    tasks = {
        i: asyncio.create_task(
            fetch_image_bytes(_build_image_query(slide.get("title", ""), topic), i)
        )
        for i, slide in enumerate(slides)
        # All slides get an image — removed `if i % 2 == 0`
    }
    images: dict[int, bytes] = {}
    for i, task in tasks.items():
        try:
            result = await task
            if result:
                images[i] = result
        except Exception as e:
            logger.warning(f"[DocumentAgent] Image task failed slide {i}: {e}")
    logger.info(f"[DocumentAgent] {len(images)}/{len(tasks)} images fetched")
    return images


# ── Agent ─────────────────────────────────────────────────────────────────────

def safe_filename(text: str, ext: str) -> str:
    slug = re.sub(r"[^\w\s-]", "", text).strip()
    slug = re.sub(r"[\s_-]+", "_", slug)[:50]
    return f"{slug}_{datetime.now().strftime('%H%M%S')}.{ext}"


class DocumentAgent(BaseAgent):
    name        = "document_agent"
    description = "Creates PowerPoint, Word, and Excel files"

    async def run(self, parameters: dict, task_id: str, context: str = "") -> Optional[str]:
        try:
            return await self._run(parameters, task_id, context)
        except asyncio.CancelledError:
            raise
        except Exception as e:
            logger.error(f"[DocumentAgent] Error: {e}", exc_info=True)
            return f"⚠️ DocumentAgent error: {e}"

    async def _run(self, parameters: dict, task_id: str, context: str) -> Optional[str]:
        all_signals = " ".join([
            parameters.get("type", ""), parameters.get("topic", ""),
            parameters.get("query", ""), parameters.get("title", ""),
            parameters.get("content", ""), parameters.get("description", ""),
            parameters.get("format", ""), context,
        ]).lower()

        PPTX_KEYS = ["presentation", "powerpoint", "pptx", "slides", "slide deck", "deck", "slideshow"]
        XLSX_KEYS = ["spreadsheet", "excel", "xlsx"]

        if any(k in all_signals for k in PPTX_KEYS):
            doc_type = "presentation"
        elif any(k in all_signals for k in XLSX_KEYS):
            doc_type = "spreadsheet"
        else:
            raw_type = parameters.get("type", "").lower()
            doc_type = (
                "presentation" if raw_type in ("presentation", "pptx", "slides")
                else "spreadsheet" if raw_type in ("spreadsheet", "excel", "xlsx")
                else "document"
            )

        topic = (
            parameters.get("topic") or parameters.get("title") or
            parameters.get("query") or parameters.get("content") or
            parameters.get("subject") or (context[:100] if context else "Untitled")
        )

        logger.info(f"[DocumentAgent] type={doc_type}  topic={topic[:60]}")

        if doc_type in ("presentation", "pptx", "slides", "powerpoint"):
            return await self._create_pptx(topic, context, parameters)
        elif doc_type in ("spreadsheet", "excel", "xlsx"):
            return await self._create_xlsx(topic, context, parameters)
        else:
            return await self._create_docx(topic, context, parameters)

    # ── PowerPoint ────────────────────────────────────────────────────────────

    async def _create_pptx(self, topic: str, context: str, params: dict) -> str:
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
        except ImportError:
            return "[DocumentAgent] python-pptx not installed. Run: pip install python-pptx"

        plan_prompt = f"""Create a professional PowerPoint presentation about: {topic}

Research context:
{context[:2000] if context else 'Use your knowledge.'}

Return ONLY valid JSON — no markdown fences:
{{
  "title": "Presentation Title",
  "subtitle": "A compelling subtitle",
  "slides": [
    {{
      "title": "Slide Title",
      "bullets": ["Point 1 with specific detail — minimum 12 words", "Point 2 with facts or examples", "Point 3", "Point 4", "Point 5"],
      "stat": "Key stat e.g. '47%' or '$2.4B' (optional)",
      "stat_label": "Short label for the stat",
      "speaker_note": "2-3 sentence speaking note for presenter"
    }}
  ]
}}

Requirements:
- Include exactly 6-8 content slides
- Each slide MUST have 5 bullet points
- Each bullet must be at least 12 words with specific facts, numbers, or examples
- Cover: overview, key concepts, real-world examples, statistics/data, challenges, future outlook
- Include a meaningful stat on at least 4 slides
- Make content deep, specific, and informative — not generic
- speaker_note should add context beyond what's on the slide"""

        raw = await self.qwen.chat(
            system_prompt="Presentation writer. Return ONLY valid JSON, no markdown fences.",
            user_message=plan_prompt,
            temperature=0.5,
        )
        raw = re.sub(r"```(?:json)?", "", raw).strip().rstrip("`").strip()
        try:
            data = json.loads(raw)
        except json.JSONDecodeError:
            logger.warning("[DocumentAgent] JSON parse failed, using fallback")
            data = {
                "title": topic, "subtitle": "Generated by AetherAI",
                "slides": [{"title": "Overview", "bullets": ["Content by AetherAI"],
                             "stat": "", "stat_label": "", "speaker_note": ""}]
            }

        T = pick_theme()
        logger.info(f"[DocumentAgent] Theme: {T['name']}")
        slides_data = data.get("slides", [])

        logger.info(f"[DocumentAgent] Fetching images for {len(slides_data)} slides…")
        slide_images = await fetch_images_for_slides(slides_data, topic)

        def rgb(tup):            return RGBColor(*tup)
        def hf():                return T["header_font"]
        def bf():                return T["body_font"]

        prs = Presentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]

        def rect(slide, x, y, w, h, color):
            s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
            s.fill.solid()
            s.fill.fore_color.rgb = color
            s.line.fill.background()
            return s

        def txt(slide, text, x, y, w, h, size, bold=False,
                color=None, align=PP_ALIGN.LEFT, italic=False, font=None):
            tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
            tf = tb.text_frame
            tf.word_wrap = True
            p  = tf.paragraphs[0]
            p.alignment = align
            r  = p.add_run()
            r.text       = str(text)
            r.font.size  = Pt(size)
            r.font.bold  = bold
            r.font.italic = italic
            r.font.name  = font or bf()
            r.font.color.rgb = color or rgb(T["text_dark"])
            return tb

        def bullets(slide, items, x, y, w, h, size=15):
            tb  = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
            tf2 = tb.text_frame
            tf2.word_wrap = True
            for j, b in enumerate(items):
                p = tf2.paragraphs[0] if j == 0 else tf2.add_paragraph()
                p.space_before = Pt(6)
                p.space_after  = Pt(5)
                r = p.add_run()
                r.text = f"▸  {b}"
                r.font.size = Pt(size)
                r.font.name = bf()
                r.font.color.rgb = rgb(T["text_dark"])

        def add_img(slide, img_bytes, x, y, w, h) -> bool:
            try:
                slide.shapes.add_picture(BytesIO(img_bytes), Inches(x), Inches(y), Inches(w), Inches(h))
                return True
            except Exception as e:
                logger.warning(f"[DocumentAgent] add_img failed: {e}")
                return False

        # ── Title slide ───────────────────────────────────────────────────────
        ts = prs.slides.add_slide(blank)
        rect(ts, 0, 0, 13.33, 7.5, rgb(T["bg_dark"]))
        rect(ts, 0, 6.2, 13.33, 1.3, rgb(T["accent"]))

        hero = await fetch_image_bytes(_build_image_query(data.get("title", topic), topic), 99, timeout=6.0)
        if hero:
            try:
                ts.shapes.add_picture(BytesIO(hero), Inches(6.5), Inches(0),
                                       Inches(6.83), Inches(7.5))
                rect(ts, 0, 0, 8.5, 7.5, rgb(T["bg_dark"]))
            except Exception:
                pass

        txt(ts, data.get("title", topic),
            0.7, 1.4, 10.5, 3.0, 44, bold=True,
            color=rgb(T["text_light"]), font=hf())
        txt(ts, data.get("subtitle", ""),
            0.7, 4.5, 9.0, 1.0, 20,
            color=rgb(T["muted"]))
        txt(ts, f"AetherAI  ·  {datetime.now().strftime('%B %d, %Y')}",
            0.7, 6.7, 11.5, 0.4, 10,
            color=rgb(T["bg_dark"]))

        # ── Content slides ────────────────────────────────────────────────────
        text_layouts = ["accent_left", "stat_card"]

        for i, sd in enumerate(slides_data):
            sl       = prs.slides.add_slide(blank)
            img      = slide_images.get(i)
            num      = str(i + 1)
            title_t  = sd.get("title", "")
            blist    = sd.get("bullets", [])
            stat     = sd.get("stat", "")
            stat_lbl = sd.get("stat_label", "")

            if img:
                # IMAGE RIGHT
                rect(sl, 0, 0, 13.33, 7.5, rgb(T["bg_light"]))
                rect(sl, 0, 0, 13.33, 1.1, rgb(T["bg_dark"]))
                txt(sl, num, 12.5, 0.18, 0.6, 0.7, 12, bold=True,
                    color=rgb(T["muted"]), align=PP_ALIGN.CENTER)
                txt(sl, title_t, 0.5, 0.13, 11.5, 0.85, 28, bold=True,
                    color=rgb(T["text_light"]), font=hf())
                bullets(sl, blist, 0.5, 1.3, 7.0, 5.8)
                add_img(sl, img, 7.8, 1.2, 5.0, 6.0)
                txt(sl, "Photo: Wikimedia / loremflickr", 7.8, 7.1, 4.5, 0.3, 8,
                    italic=True, color=rgb(T["muted"]))

            else:
                layout = text_layouts[i % len(text_layouts)]

                if layout == "accent_left":
                    # ACCENT LEFT
                    rect(sl, 0, 0, 3.8, 7.5, rgb(T["bg_dark"]))
                    rect(sl, 3.8, 0, 9.53, 7.5, rgb(T["bg_light"]))
                    txt(sl, num, 0.3, 0.3, 3.0, 0.8, 32, bold=True,
                        color=rgb(T["accent"]))
                    txt(sl, title_t, 0.3, 1.2, 3.2, 4.0, 20, bold=True,
                        color=rgb(T["text_light"]), font=hf())
                    bullets(sl, blist, 4.1, 0.7, 8.8, 6.0)
                    if stat:
                        txt(sl, stat, 0.3, 5.7, 3.2, 1.0, 28, bold=True,
                            color=rgb(T["accent"]))
                        txt(sl, stat_lbl, 0.3, 6.5, 3.2, 0.6, 10,
                            color=rgb(T["muted"]))

                else:
                    # STAT CARD
                    rect(sl, 0, 0, 13.33, 7.5, rgb(T["bg_light"]))
                    rect(sl, 0, 0, 13.33, 1.9, rgb(T["accent"]))
                    txt(sl, num, 12.4, 0.2, 0.7, 0.6, 11, bold=True,
                        color=rgb(T["bg_dark"]), align=PP_ALIGN.CENTER)
                    txt(sl, title_t, 0.5, 0.25, 11.5, 1.4, 32, bold=True,
                        color=rgb(T["text_light"]), font=hf())
                    if stat:
                        rect(sl, 0.5, 2.1, 3.5, 4.5, rgb(T["bg_dark"]))
                        txt(sl, stat, 0.6, 2.5, 3.2, 2.0, 52, bold=True,
                            color=rgb(T["accent"]), align=PP_ALIGN.CENTER)
                        txt(sl, stat_lbl, 0.6, 4.6, 3.2, 0.8, 12,
                            color=rgb(T["muted"]), align=PP_ALIGN.CENTER)
                        bullets(sl, blist, 4.3, 2.0, 8.5, 5.0)
                    else:
                        bullets(sl, blist, 0.5, 2.1, 12.3, 5.0)

            note = sd.get("speaker_note", "")
            if note:
                sl.notes_slide.notes_text_frame.text = note

        # ── Closing slide ─────────────────────────────────────────────────────
        cs = prs.slides.add_slide(blank)
        rect(cs, 0, 0, 13.33, 7.5, rgb(T["bg_dark"]))
        rect(cs, 0, 0, 13.33, 3.0, rgb(T["accent"]))
        txt(cs, "Thank You", 0.8, 0.4, 11.5, 2.0, 60, bold=True,
            color=rgb(T["bg_dark"]), align=PP_ALIGN.CENTER, font=hf())
        txt(cs, data.get("title", topic), 0.8, 3.5, 11.5, 1.0, 18,
            color=rgb(T["muted"]), align=PP_ALIGN.CENTER)
        txt(cs, "Generated by AetherAI", 0.8, 4.6, 11.5, 0.5, 11,
            color=rgb(T["muted"]), align=PP_ALIGN.CENTER)

        fname = safe_filename(data.get("title", topic), "pptx")
        fpath = OUTPUT_DIR / fname
        prs.save(str(fpath))

        n = len(slides_data)
        return (
            f"✅ PowerPoint created: output/{fname}\n"
            f"Theme: {T['name']}  |  Slides: {n + 2}  |  Photos: {len(slide_images)}\n"
            f"Topic: {data.get('title', topic)}\n"
            f"Full path: {fpath}"
        )

    # ── Word Document ─────────────────────────────────────────────────────────

    async def _create_docx(self, topic: str, context: str, params: dict) -> str:
        try:
            from docx import Document
            from docx.shared import Pt, Inches, RGBColor as DR
            from docx.enum.text import WD_ALIGN_PARAGRAPH
        except ImportError:
            return "[DocumentAgent] python-docx not installed. Run: pip install python-docx"

        def _extract_json(raw: str):
            """Try multiple strategies to extract valid JSON from LLM output."""
            import json as _j, re as _r
            cleaned = _r.sub(r"```(?:json)?", "", raw).strip().rstrip("`").strip()
            try: return _j.loads(cleaned)
            except Exception: pass
            s, e = cleaned.find("{"), cleaned.rfind("}")
            if s != -1 and e > s:
                try: return _j.loads(cleaned[s:e+1])
                except Exception: pass
            return None

        prompt = f"""Write a professional, detailed document about: {topic}

Context:
{context[:2000] if context else 'Use your knowledge.'}

Return ONLY valid JSON — no markdown fences, no extra text:
{{
  "title": "Document Title",
  "sections": [
    {{"heading": "Section Heading", "paragraphs": ["Paragraph 1 — at least 60 words with specific facts and detail...", "Paragraph 2 with supporting information and examples..."]}}
  ],
  "conclusion": "Concluding paragraph of at least 60 words with key takeaways and recommendations."
}}

Requirements:
- Include exactly 5-6 sections covering: introduction, key background, main analysis, real-world examples, challenges/limitations, future outlook
- Each section must have 2-3 paragraphs
- Each paragraph must be at least 60 words with specific facts, numbers, and concrete details
- No generic filler — be thorough and informative about the topic"""

        raw = await self.qwen.chat(
            system_prompt="Professional document writer. Return ONLY valid JSON. No markdown fences. No text before or after the JSON object.",
            user_message=prompt, temperature=0.5,
        )
        data = _extract_json(raw)

        # ── Retry with simplified prompt if JSON extraction failed ────────────
        if not data or not data.get("sections"):
            logger.warning(f"[DocumentAgent] DOCX JSON parse failed, retrying...")
            retry_prompt = f"""Write about: {topic}

Return ONLY this JSON (fill in the placeholder text with real content):
{{"title": "{topic.title()}", "sections": [{{"heading": "Introduction", "paragraphs": ["WRITE 80+ WORD INTRODUCTION HERE"]}}, {{"heading": "Background", "paragraphs": ["WRITE 80+ WORD BACKGROUND HERE"]}}, {{"heading": "Key Analysis", "paragraphs": ["WRITE 80+ WORD ANALYSIS HERE"]}}, {{"heading": "Examples and Applications", "paragraphs": ["WRITE 80+ WORD EXAMPLES HERE"]}}, {{"heading": "Challenges and Future Outlook", "paragraphs": ["WRITE 80+ WORD CHALLENGES HERE"]}}], "conclusion": "WRITE 80+ WORD CONCLUSION HERE"}}

Replace every placeholder like "WRITE 80+ WORD ... HERE" with real paragraph text about {topic}."""

            raw2 = await self.qwen.chat(
                system_prompt="Return ONLY valid JSON. Replace all placeholder text with real content. No markdown.",
                user_message=retry_prompt, temperature=0.3,
            )
            data = _extract_json(raw2)

        # ── Final fallback: generate prose and wrap in structure ──────────────
        if not data or not data.get("sections"):
            logger.warning(f"[DocumentAgent] Both JSON attempts failed, using prose fallback")
            generated = await self.qwen.chat(
                system_prompt="You are a professional writer. Write a detailed, well-structured document.",
                user_message=f"Write a comprehensive document about: {topic}\n\nInclude a detailed introduction, 3-4 main sections with thorough paragraphs (at least 80 words each), and a conclusion. Cover the topic in depth.",
                temperature=0.6,
            )
            # Split generated text into sections by detecting headings
            sections = []
            current_heading = "Overview"
            current_paras = []
            for line in generated.split("\n"):
                line = line.strip()
                if not line:
                    continue
                # Detect heading-like lines (short, ends with no period, possibly bold)
                is_heading = (len(line) < 60 and not line.endswith(".") and
                              (line.startswith("#") or line.isupper() or
                               (len(line.split()) <= 6 and line[0].isupper())))
                if is_heading and current_paras:
                    sections.append({"heading": current_heading, "paragraphs": current_paras})
                    current_heading = line.lstrip("#").strip()
                    current_paras = []
                elif not is_heading:
                    current_paras.append(line)
            if current_paras:
                sections.append({"heading": current_heading, "paragraphs": current_paras})
            if not sections:
                sections = [{"heading": "Content", "paragraphs": [generated[:3000]]}]
            data = {"title": topic.title(), "sections": sections, "conclusion": ""}

        T = pick_theme()

        def dr(tup): return DR(*tup)

        doc = Document()
        for sec in doc.sections:
            sec.top_margin = sec.bottom_margin = Inches(1.0)
            sec.left_margin = sec.right_margin = Inches(1.25)

        tp = doc.add_paragraph()
        tp.alignment = WD_ALIGN_PARAGRAPH.CENTER
        tr = tp.add_run(data.get("title", topic))
        tr.bold = True
        tr.font.size = Pt(26)
        tr.font.name = T["header_font"]
        tr.font.color.rgb = dr(T["bg_dark"])

        dp = doc.add_paragraph()
        dp.alignment = WD_ALIGN_PARAGRAPH.CENTER
        dr2 = dp.add_run(
            f"{T['name']} Theme  ·  Generated by AetherAI  ·  "
            f"{datetime.now().strftime('%B %d, %Y')}"
        )
        dr2.font.size = Pt(10)
        dr2.font.color.rgb = dr(T["muted"])
        doc.add_paragraph()

        for sec in data.get("sections", []):
            h = doc.add_heading(sec.get("heading", ""), level=1)
            if h.runs:
                h.runs[0].font.color.rgb = dr(T["bg_dark"])
                h.runs[0].font.size = Pt(14)
                h.runs[0].font.name = T["header_font"]
            for pt in sec.get("paragraphs", []):
                p = doc.add_paragraph(pt)
                p.paragraph_format.space_after = Pt(8)
                p.paragraph_format.first_line_indent = Inches(0.3)

        if data.get("conclusion"):
            h2 = doc.add_heading("Conclusion", level=1)
            if h2.runs:
                h2.runs[0].font.color.rgb = dr(T["bg_dark"])
                h2.runs[0].font.size = Pt(14)
            doc.add_paragraph(data["conclusion"])

        fname = safe_filename(data.get("title", topic), "docx")
        fpath = OUTPUT_DIR / fname
        doc.save(str(fpath))
        return (
            f"✅ Word document created: output/{fname}\n"
            f"Theme: {T['name']}  |  Sections: {len(data.get('sections', []))}\n"
            f"Title: {data.get('title', topic)}\nFull path: {fpath}"
        )
    
    
    # ── Excel Spreadsheet ─────────────────────────────────────────────────────

    async def _create_xlsx(self, topic: str, context: str, params: dict) -> str:
        try:
            import openpyxl
            from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
            from openpyxl.utils import get_column_letter
        except ImportError:
            return "[DocumentAgent] openpyxl not installed. Run: pip install openpyxl"

        prompt = f"""Create comprehensive spreadsheet data about: {topic}

Context: {context[:1500] if context else 'Use your knowledge.'}

Return ONLY valid JSON:
{{
  "title": "Spreadsheet Title",
  "sheets": [
    {{
      "name": "Sheet Name",
      "headers": ["Col1", "Col2", "Col3", "Col4", "Col5"],
      "rows": [["val1","val2",100,45.5,"note"]],
      "has_totals": true
    }}
  ]
}}

Requirements:
- Include 2-3 sheets covering different aspects of "{topic}"
- Each sheet must have 5-7 columns and at least 12 data rows
- Include numeric columns (integers or decimals) where appropriate
- Make data realistic, varied, and specific to "{topic}"
- has_totals: true adds a SUM row at the bottom for numeric columns"""

        raw = await self.qwen.chat(
            system_prompt="Data analyst. Return ONLY valid JSON, no markdown fences.",
            user_message=prompt, temperature=0.4,
        )
        raw = re.sub(r"```(?:json)?", "", raw).strip().rstrip("`").strip()
        try:
            data = json.loads(raw)
        except json.JSONDecodeError:
            data = {"title": topic,
                    "sheets": [{"name": "Data", "headers": ["Item", "Value"],
                                 "rows": [["Example", "Data"]]}]}

        T = pick_theme()

        def hx(tup): return "{:02X}{:02X}{:02X}".format(*tup)

        wb = openpyxl.Workbook()
        wb.remove(wb.active)

        hdr_fill   = PatternFill("solid", fgColor=hx(T["bg_dark"]))
        alt_fill   = PatternFill("solid", fgColor=hx(T["alt_row"]))
        title_font = Font(bold=True, color=hx(T["bg_dark"]),
                          size=15, name=T["header_font"])
        hdr_font   = Font(bold=True, color=hx(T["text_light"]),
                          size=11, name=T["body_font"])
        thin   = Side(style="thin", color="CCCCCC")
        border = Border(left=thin, right=thin, top=thin, bottom=thin)
        center = Alignment(horizontal="center", vertical="center")
        left   = Alignment(horizontal="left",   vertical="center")

        for sh in data.get("sheets", []):
            ws      = wb.create_sheet(title=sh.get("name", "Sheet")[:31])
            headers = sh.get("headers", [])
            ncols   = max(len(headers), 1)

            ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=ncols)
            c = ws.cell(1, 1, data.get("title", topic))
            c.font = title_font
            c.alignment = center
            c.fill = PatternFill("solid", fgColor=hx(T["bg_light"]))
            ws.row_dimensions[1].height = 32

            for ci in range(1, ncols + 1):
                ws.cell(2, ci, "").fill = PatternFill("solid", fgColor=hx(T["accent"]))
            ws.row_dimensions[2].height = 4

            for ci, h in enumerate(headers, 1):
                cell = ws.cell(3, ci, h)
                cell.font = hdr_font
                cell.fill = hdr_fill
                cell.alignment = center
                cell.border = border
            ws.row_dimensions[3].height = 24

            for ri, row in enumerate(sh.get("rows", []), 4):
                for ci, val in enumerate(row, 1):
                    cell = ws.cell(ri, ci, val)
                    cell.border = border
                    cell.alignment = left
                    if ri % 2 == 0:
                        cell.fill = alt_fill
                ws.row_dimensions[ri].height = 18

            # Totals row for numeric columns
            rows_data = sh.get("rows", [])
            if sh.get("has_totals") and rows_data:
                totals_row = 4 + len(rows_data)
                total_fill = PatternFill("solid", fgColor=hx(T["bg_dark"]))
                for ci, h in enumerate(headers, 1):
                    cell = ws.cell(totals_row, ci)
                    cell.fill = total_fill
                    cell.border = border
                    # Check if column has numeric values
                    col_vals = [r[ci-1] for r in rows_data if ci-1 < len(r)]
                    if all(isinstance(v, (int, float)) for v in col_vals if v != "" and v is not None):
                        from openpyxl.utils import get_column_letter as gcl
                        col_l = gcl(ci)
                        cell.value = f"=SUM({col_l}4:{col_l}{totals_row-1})"
                        cell.font = Font(bold=True, color=hx(T["text_light"]), size=11, name=T["body_font"])
                        cell.alignment = center
                    elif ci == 1:
                        cell.value = "TOTAL"
                        cell.font = Font(bold=True, color=hx(T["text_light"]), size=11, name=T["body_font"])
                        cell.alignment = left
                    else:
                        cell.font = Font(bold=True, color=hx(T["text_light"]), size=11)
                ws.row_dimensions[totals_row].height = 22

            for ci, _ in enumerate(headers, 1):
                col_letter = get_column_letter(ci)
                max_len    = len(str(headers[ci - 1]))
                for row in sh.get("rows", []):
                    if ci <= len(row):
                        max_len = max(max_len, len(str(row[ci - 1])))
                ws.column_dimensions[col_letter].width = min(max_len + 4, 40)

        fname = safe_filename(data.get("title", topic), "xlsx")
        fpath = OUTPUT_DIR / fname
        wb.save(str(fpath))
        sheets_info = ", ".join(s.get("name", "Sheet") for s in data.get("sheets", []))
        return (
            f"✅ Excel spreadsheet created: output/{fname}\n"
            f"Theme: {T['name']}  |  Sheets: {sheets_info}\n"
            f"Title: {data.get('title', topic)}\nFull path: {fpath}"
        )
